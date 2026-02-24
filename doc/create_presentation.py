"""Generate MCP Gateway workshop PowerPoint presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Colors (AI Devs inspired dark theme) ────────────────────────
BG_DARK = RGBColor(0x0D, 0x11, 0x17)
BG_CARD = RGBColor(0x16, 0x1B, 0x22)
BG_CARD_LIGHT = RGBColor(0x21, 0x26, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x8B, 0x94, 0x9E)
LIGHT_GRAY = RGBColor(0xC9, 0xD1, 0xD9)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
RED = RGBColor(0xEF, 0x44, 0x44)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
PINK = RGBColor(0xEC, 0x48, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helper functions ────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txbox


def add_para(tf, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = align
    p.space_before = space_before
    return p


def add_box(slide, left, top, width, height, fill_color=BG_CARD, border_color=None, radius=Inches(0.15)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top),
                                   Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=CYAN, width=Pt(3)):
    connector = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x1), Inches(y1),
        Inches(x2 - x1), Inches(0.35)
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = color
    connector.line.fill.background()
    return connector


def add_line_arrow(slide, x1, y1, x2, y2, color=CYAN, width=Pt(2)):
    """Add a line with arrow using a connector."""
    from pptx.oxml.ns import qn
    connector = slide.shapes.add_connector(
        1,  # straight connector
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead
    ln = connector.line._ln
    tail_end = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail_end)
    return connector


def add_labeled_box(slide, text, left, top, width, height, fill_color=BG_CARD,
                    border_color=CYAN, text_color=WHITE, font_size=14, bold=True):
    box = add_box(slide, left, top, width, height, fill_color, border_color)
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(0)
    # Vertical center
    from pptx.enum.text import MSO_ANCHOR
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return box


def add_icon_circle(slide, icon_text, left, top, size, bg_color, text_color=WHITE, font_size=20):
    circle = add_circle(slide, left, top, size, bg_color)
    tf = circle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    return circle


def add_subtitle_line(slide):
    """Add a thin accent line under the title area."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45),
                                   Inches(2), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.fill.background()


# ════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════

slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1)

# Decorative gradient circles (top-right, bottom-left)
c1 = add_circle(slide1, 9.5, -1.5, 5, PURPLE)
c1.fill.solid()
c1.fill.fore_color.rgb = PURPLE
from pptx.oxml.ns import qn
# Make semi-transparent
c1_sp = c1._element
solid = c1_sp.find(qn('p:spPr')).find(qn('a:solidFill'))
if solid is not None:
    srgb = solid.find(qn('a:srgbClr'))
    if srgb is not None:
        alpha = srgb.makeelement(qn('a:alpha'), {'val': '15000'})
        srgb.append(alpha)

c2 = add_circle(slide1, -2, 4, 5, CYAN)
c2_sp = c2._element
solid2 = c2_sp.find(qn('p:spPr')).find(qn('a:solidFill'))
if solid2 is not None:
    srgb2 = solid2.find(qn('a:srgbClr'))
    if srgb2 is not None:
        alpha2 = srgb2.makeelement(qn('a:alpha'), {'val': '10000'})
        srgb2.append(alpha2)

# Title
add_text(slide1, "MCP", 1.5, 1.8, 10, 1.2, size=72, color=CYAN, bold=True, align=PP_ALIGN.LEFT)
add_text(slide1, "w użyciu korporacyjnym", 1.5, 3.0, 10, 0.8, size=40, color=WHITE, bold=True, align=PP_ALIGN.LEFT)

# Separator line
shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.0), Inches(3), Inches(0.05))
shape.fill.solid()
shape.fill.fore_color.rgb = CYAN
shape.line.fill.background()

# Presenter
add_text(slide1, "Adam Dąbrowski", 1.5, 4.3, 6, 0.6, size=24, color=LIGHT_GRAY, bold=False)
add_text(slide1, "Warsztaty  |  Model Context Protocol", 1.5, 4.9, 8, 0.5, size=16, color=GRAY)

# Decorative MCP protocol icons on right side
for i, (txt, clr, yy) in enumerate([
    ("{ }", CYAN, 2.5), ("< >", PURPLE, 3.5), ("[ ]", GREEN, 4.5)
]):
    add_text(slide1, txt, 10.5, yy, 2, 0.6, size=36, color=clr, bold=True, align=PP_ALIGN.CENTER,
             font_name="Consolas")


# ════════════════════════════════════════════════════════════════
# SLIDE 2: What is an Agent
# ════════════════════════════════════════════════════════════════

slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)

add_text(slide2, "Czym jest Agent?", 0.8, 0.4, 8, 0.8, size=36, color=CYAN, bold=True)
add_subtitle_line(slide2)

# ── Main diagram: User → Agent (loop) → LLM / Tools ──

# User box
add_labeled_box(slide2, "Użytkownik", 0.5, 2.3, 1.6, 0.9, BG_CARD, PURPLE, WHITE, 14)

# Arrow User → Agent
add_arrow(slide2, 2.2, 2.5, 3.0, 2.5, PURPLE)

# Agent box (larger, with loop indicator)
agent_box = add_box(slide2, 3.1, 1.8, 2.8, 2.2, BG_CARD, CYAN)
tf = agent_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Agent"
p.font.size = Pt(20)
p.font.color.rgb = CYAN
p.font.bold = True
p.font.name = "Segoe UI"
p.alignment = PP_ALIGN.CENTER
add_para(tf, "Agentic Loop", size=11, color=GRAY, align=PP_ALIGN.CENTER)
add_para(tf, "↻ Reason → Act → Observe", size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Loop arrow (circular indicator on agent)
loop_indicator = add_text(slide2, "↻", 5.4, 1.6, 0.6, 0.5, size=28, color=CYAN, bold=True)

# Arrow Agent → LLM
add_arrow(slide2, 6.0, 2.1, 6.8, 2.1, CYAN)

# LLM box
add_labeled_box(slide2, "LLM", 6.9, 1.8, 1.5, 0.9, BG_CARD, CYAN, WHITE, 16)

# Arrow Agent → Tools Internal
add_arrow(slide2, 6.0, 2.8, 6.8, 2.8, GREEN)

# Internal Tools box
add_labeled_box(slide2, "Tools\n(wewnętrzne)", 6.9, 2.9, 1.5, 0.9, BG_CARD, GREEN, WHITE, 12)

# Arrow Agent → MCP Tools External
add_arrow(slide2, 6.0, 3.5, 6.8, 3.5, PURPLE)

# External Tools box
ext_box = add_labeled_box(slide2, "Tools\n(zewnętrzne)", 6.9, 3.9, 1.5, 0.9, BG_CARD, PURPLE, WHITE, 12)

# MCP label on external tools
add_text(slide2, "MCP", 8.5, 4.0, 1.0, 0.4, size=14, color=PURPLE, bold=True)

# ── Bottom: Agent type matrix ──

add_text(slide2, "Rodzaje agentów", 0.5, 5.2, 4, 0.5, size=18, color=LIGHT_GRAY, bold=True)

# Matrix grid
matrix_x = 0.8
matrix_y = 5.8
cell_w = 2.8
cell_h = 0.7

# Headers
add_text(slide2, "", matrix_x, matrix_y - 0.05, cell_w, 0.3, size=11, color=GRAY)
add_text(slide2, "W imieniu użytkownika", matrix_x + cell_w + 0.1, matrix_y - 0.35, cell_w, 0.3,
         size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide2, "Autonomiczny", matrix_x + 2 * (cell_w + 0.1), matrix_y - 0.35, cell_w, 0.3,
         size=12, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)

# Row labels
add_text(slide2, "Local", matrix_x - 0.1, matrix_y + 0.1, 1.0, cell_h,
         size=13, color=GREEN, bold=True, align=PP_ALIGN.RIGHT)
add_text(slide2, "Remote", matrix_x - 0.1, matrix_y + cell_h + 0.2, 1.0, cell_h,
         size=13, color=ORANGE, bold=True, align=PP_ALIGN.RIGHT)

# Matrix cells
cells = [
    ("Claude Code\nCursor, Copilot", matrix_x + cell_w + 0.1, matrix_y, CYAN),
    ("Local cron agents\nFile watchers", matrix_x + 2 * (cell_w + 0.1), matrix_y, PURPLE),
    ("Chatbot webowy\nAPI assistant", matrix_x + cell_w + 0.1, matrix_y + cell_h + 0.2, CYAN),
    ("Monitoring agent\nPipeline agent", matrix_x + 2 * (cell_w + 0.1), matrix_y + cell_h + 0.2, PURPLE),
]
for text, cx, cy, border in cells:
    box = add_box(slide2, cx, cy, cell_w, cell_h, BG_CARD, border)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════════
# SLIDE 3: What is MCP
# ════════════════════════════════════════════════════════════════

slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)

add_text(slide3, "Model Context Protocol (MCP)", 0.8, 0.4, 10, 0.8, size=36, color=CYAN, bold=True)
add_subtitle_line(slide3)

# Left side: description
add_text(slide3, "Otwarty standard komunikacji\nmiędzy agentami AI a narzędziami", 0.8, 1.7, 5.5, 1.0,
         size=22, color=WHITE, bold=True)

desc_box = add_text(slide3, "", 0.8, 2.8, 5.2, 3.5, size=16, color=LIGHT_GRAY)
tf = desc_box.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = ""

bullets = [
    ("Jeden protokół", "zamiast N integracji"),
    ("JSON-RPC 2.0", "przez HTTP (Streamable HTTP)"),
    ("Standaryzacja", "tools, resources, prompts"),
    ("Open source", "specyfikacja Anthropic (2024)"),
]
for title, desc in bullets:
    p = tf.add_paragraph()
    p.space_before = Pt(12)
    run1 = p.add_run()
    run1.text = f"▸ {title}  "
    run1.font.size = Pt(17)
    run1.font.color.rgb = CYAN
    run1.font.bold = True
    run1.font.name = "Segoe UI"
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(15)
    run2.font.color.rgb = LIGHT_GRAY
    run2.font.name = "Segoe UI"

# Right side: MCP diagram
# Host/Client/Server model
diagram_x = 7.0

# Host
host_box = add_box(slide3, diagram_x, 1.7, 5.5, 5.0, BG_CARD, GRAY)
add_text(slide3, "MCP Host", diagram_x + 0.2, 1.8, 2, 0.4, size=13, color=GRAY, bold=True)

# Client inside host
client_box = add_box(slide3, diagram_x + 0.3, 2.4, 2.2, 1.5, BG_CARD_LIGHT, CYAN)
add_text(slide3, "MCP Client", diagram_x + 0.5, 2.5, 2, 0.35, size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide3, "Maintains 1:1\nconnection", diagram_x + 0.5, 2.9, 2, 0.6, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# Second client
client_box2 = add_box(slide3, diagram_x + 0.3, 4.2, 2.2, 1.5, BG_CARD_LIGHT, CYAN)
add_text(slide3, "MCP Client", diagram_x + 0.5, 4.3, 2, 0.35, size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide3, "Maintains 1:1\nconnection", diagram_x + 0.5, 4.7, 2, 0.6, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# Arrows to servers
add_arrow(slide3, diagram_x + 2.6, 3.0, diagram_x + 3.2, 3.0, CYAN, Pt(2))
add_arrow(slide3, diagram_x + 2.6, 4.8, diagram_x + 3.2, 4.8, CYAN, Pt(2))

# MCP Servers
server_box1 = add_box(slide3, diagram_x + 3.3, 2.4, 1.8, 1.5, BG_CARD_LIGHT, GREEN)
add_text(slide3, "MCP Server", diagram_x + 3.4, 2.5, 1.6, 0.35, size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide3, "Tools\nResources\nPrompts", diagram_x + 3.4, 2.9, 1.6, 0.9, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

server_box2 = add_box(slide3, diagram_x + 3.3, 4.2, 1.8, 1.5, BG_CARD_LIGHT, GREEN)
add_text(slide3, "MCP Server", diagram_x + 3.4, 4.3, 1.6, 0.35, size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide3, "Tools\nResources\nPrompts", diagram_x + 3.4, 4.7, 1.6, 0.9, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 4: Organizational Spaghetti
# ════════════════════════════════════════════════════════════════

slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)

add_text(slide4, "Spaghetti MCP", 0.8, 0.4, 10, 0.8, size=36, color=RED, bold=True)

shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45),
                                Inches(2), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = RED
shape.line.fill.background()

add_text(slide4, "Każdy agent podłączony do każdego serwera MCP", 0.8, 1.7, 10, 0.5,
         size=18, color=GRAY)

# ── Left column: Agents ──
agents = [
    ("Agent 1\n(Sales)", 0.5, 2.5, CYAN),
    ("Agent 2\n(Support)", 0.5, 3.5, CYAN),
    ("Agent 3\n(DevOps)", 0.5, 4.5, CYAN),
    ("Agent 4\n(HR)", 0.5, 5.5, CYAN),
    ("Agent 5\n(Finance)", 0.5, 6.3, CYAN),
]

# ── Right column: MCP Servers ──
mcp_servers = [
    ("MCP: CRM", 9.5, 2.3, TEAL),
    ("MCP: Jira", 9.5, 3.0, TEAL),
    ("MCP: Git", 9.5, 3.7, TEAL),
    ("MCP: Confluence", 9.5, 4.4, TEAL),
    ("MCP: Slack", 9.5, 5.1, TEAL),
    ("MCP: HR System", 9.5, 5.8, TEAL),
    ("MCP: DB", 9.5, 6.5, TEAL),
]

# Draw agents
for name, x, y, clr in agents:
    add_labeled_box(slide4, name, x, y, 1.8, 0.7, BG_CARD, clr, WHITE, 11)

# Draw MCP servers
for name, x, y, clr in mcp_servers:
    add_labeled_box(slide4, name, x, y, 2.2, 0.55, BG_CARD, clr, WHITE, 11)

# Draw spaghetti lines — every agent to every MCP server
for _, ax, ay, _ in agents:
    agent_cx = ax + 1.8  # right edge of agent box
    agent_cy = ay + 0.35  # vertical center
    for _, sx, sy, _ in mcp_servers:
        server_cy = sy + 0.275
        try:
            conn = slide4.shapes.add_connector(1,
                Inches(agent_cx), Inches(agent_cy),
                Inches(sx), Inches(server_cy))
            conn.line.color.rgb = RGBColor(0xFF, 0x6B, 0x35)
            conn.line.width = Pt(1)
            ln = conn.line._ln
            solid_fill = ln.find(qn('a:solidFill'))
            if solid_fill is not None:
                srgb = solid_fill.find(qn('a:srgbClr'))
                if srgb is not None:
                    alpha = srgb.makeelement(qn('a:alpha'), {'val': '30000'})
                    srgb.append(alpha)
        except Exception:
            pass

# Center labels highlighting problems
add_labeled_box(slide4, "N x M\npołączeń", 5.0, 2.8, 1.8, 0.8, BG_CARD_LIGHT, RED, RED, 14)
add_labeled_box(slide4, "N x M\ntokenów", 5.0, 3.9, 1.8, 0.8, BG_CARD_LIGHT, ORANGE, ORANGE, 14)
add_labeled_box(slide4, "N x M\nservice\naccounts", 5.0, 5.0, 1.8, 1.0, BG_CARD_LIGHT, YELLOW, YELLOW, 13)

# Bottom warning
add_text(slide4, "⚠  5 agentów × 7 serwerów MCP = 35 połączeń, 35 tokenów, zero governance",
         0.8, 7.0, 12, 0.3, size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 5: Problems with MCP
# ════════════════════════════════════════════════════════════════

slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)

add_text(slide5, "Problemy z MCP w korporacji", 0.8, 0.4, 10, 0.8, size=36, color=ORANGE, bold=True)
shape = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45),
                                Inches(2), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

# Three problem cards
problems = [
    {
        "title": "Discovery",
        "icon": "🔍",
        "color": ORANGE,
        "desc": "Skąd agent wie jakie serwery MCP\nsą dostępne w organizacji?",
        "details": [
            "Brak centralnego rejestru",
            "Hardkodowane URL-e w konfiguracji",
            "Każdy zespół ma swoje serwery",
            "Agent musi znać wszystkie z góry",
        ]
    },
    {
        "title": "Context Rotting",
        "icon": "🧠",
        "color": YELLOW,
        "desc": "Wszystkie serwery MCP ładowane\nna start — przeładowany kontekst",
        "details": [
            "Wszystkie tools w kontekście LLM",
            "Agent nie wie co wybrać",
            "Halucynacje przy wyborze toola",
            "Nie skaluje się (50+ serwerów)",
        ]
    },
    {
        "title": "Autoryzacja",
        "icon": "🔒",
        "color": RED,
        "desc": "Każdy serwer MCP wymaga\nosobnej autoryzacji",
        "details": [
            "N serwerów = N tokenów",
            "Brak SSO dla agentów",
            "Service accounts bez kontroli",
            "Brak audytu per-user",
        ]
    },
]

card_w = 3.6
card_h = 4.6
start_x = 0.8
gap = 0.5

for i, prob in enumerate(problems):
    x = start_x + i * (card_w + gap)
    y = 2.0

    # Card background
    card = add_box(slide5, x, y, card_w, card_h, BG_CARD, prob["color"])

    # Icon + Title
    add_text(slide5, prob["icon"], x + 0.2, y + 0.2, 0.6, 0.6, size=28, color=prob["color"])
    add_text(slide5, prob["title"], x + 0.8, y + 0.25, card_w - 1, 0.5,
             size=22, color=prob["color"], bold=True)

    # Description
    add_text(slide5, prob["desc"], x + 0.3, y + 0.9, card_w - 0.6, 0.8,
             size=13, color=LIGHT_GRAY)

    # Separator
    sep = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x + 0.3), Inches(y + 1.8),
                                  Inches(card_w - 0.6), Inches(0.02))
    sep.fill.solid()
    sep.fill.fore_color.rgb = prob["color"]
    sep.line.fill.background()

    # Details
    for j, detail in enumerate(prob["details"]):
        add_text(slide5, f"▸  {detail}", x + 0.3, y + 2.0 + j * 0.55, card_w - 0.6, 0.5,
                 size=13, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 6: MCP Gateway as Solution (Concept)
# ════════════════════════════════════════════════════════════════

slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6)

add_text(slide6, "MCP Gateway — rozwiązanie", 0.8, 0.4, 10, 0.8, size=36, color=GREEN, bold=True)
shape = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45),
                                Inches(2), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = GREEN
shape.line.fill.background()

# ── Left: Agents ──
add_text(slide6, "Agenci", 0.3, 1.8, 2, 0.4, size=14, color=GRAY, bold=True, align=PP_ALIGN.CENTER)

agent_items = [
    ("Agent\n(Remote)", 0.5, 2.3, CYAN),
    ("Agent\n(Local)", 0.5, 3.5, GREEN),
    ("Agent\n(Auto)", 0.5, 4.7, PURPLE),
]
for name, ax, ay, clr in agent_items:
    add_labeled_box(slide6, name, ax, ay, 1.6, 0.9, BG_CARD, clr, WHITE, 12)

# ── Center-left: IDP ──
# IDP box
idp_box = add_box(slide6, 2.8, 2.0, 1.8, 4.0, BG_CARD, YELLOW)
add_text(slide6, "IDP", 2.9, 2.1, 1.6, 0.4, size=18, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)
add_text(slide6, "Keycloak", 2.9, 2.55, 1.6, 0.3, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# OIDC Login
add_text(slide6, "OIDC\nLogin", 3.0, 3.0, 1.4, 0.7, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Token Exchange
te_box = add_box(slide6, 3.0, 3.8, 1.4, 0.9, BG_CARD_LIGHT, YELLOW)
add_text(slide6, "Token\nExchange", 3.05, 3.85, 1.3, 0.7, size=11, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)

# JWKS
add_text(slide6, "JWKS\nEndpoint", 3.0, 4.9, 1.4, 0.7, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Arrows agents → IDP
for ay in [2.55, 3.75, 4.95]:
    add_arrow(slide6, 2.15, ay, 2.7, ay, YELLOW, Pt(2))

# ── Center: Gateway ──
gw_x = 5.2
gw_box = add_box(slide6, gw_x, 1.8, 3.0, 5.0, BG_CARD, GREEN)
add_text(slide6, "MCP Gateway", gw_x + 0.1, 1.9, 2.8, 0.45, size=20, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# Features inside gateway
gw_features = [
    ("JWT Validation", 2.5, CYAN),
    ("Token Exchange", 3.2, YELLOW),
    ("Discovery", 3.9, GREEN),
    ("Lazy Init", 4.6, TEAL),
    ("Tool Proxy", 5.3, PURPLE),
]
for feat, fy, fc in gw_features:
    feat_box = add_box(slide6, gw_x + 0.2, fy, 2.6, 0.55, BG_CARD_LIGHT, fc)
    add_text(slide6, feat, gw_x + 0.3, fy + 0.05, 2.4, 0.4, size=12, color=fc, bold=True, align=PP_ALIGN.CENTER)

# Arrows IDP ↔ Gateway
add_arrow(slide6, 4.65, 3.1, 5.1, 3.1, YELLOW, Pt(2))
add_arrow(slide6, 4.65, 4.0, 5.1, 4.0, YELLOW, Pt(2))

# ── Right: MCP Servers ──
add_text(slide6, "MCP Servers", 9.0, 1.8, 3, 0.4, size=14, color=GRAY, bold=True, align=PP_ALIGN.CENTER)

servers = [
    ("Weather\nServer", 9.0, 2.4, TEAL, "aud: mcp-weather"),
    ("Calculator\nServer", 9.0, 3.6, TEAL, "aud: mcp-calc"),
    ("CRM\nServer", 9.0, 4.8, TEAL, "aud: mcp-crm"),
    ("...", 9.0, 6.0, GRAY, ""),
]
for name, sx, sy, sc, aud in servers:
    add_labeled_box(slide6, name, sx, sy, 1.8, 0.8, BG_CARD, sc, WHITE, 12)
    if aud:
        add_text(slide6, aud, 10.9, sy + 0.15, 2, 0.4, size=9, color=GRAY,
                 font_name="Consolas")

# JWT offline verification note
add_text(slide6, "JWT offline\nverification", 11.0, 5.3, 1.8, 0.6, size=10, color=YELLOW, bold=True)

# Arrows Gateway → Servers
for sy in [2.6, 3.8, 5.0]:
    add_arrow(slide6, 8.25, sy, 8.9, sy, TEAL, Pt(2))

# Bottom label
add_text(slide6, "1 login  →  1 token  →  N serwerów (z automatycznym Token Exchange)",
         0.8, 6.9, 12, 0.4, size=15, color=GREEN, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 7: DEMO Architecture
# ════════════════════════════════════════════════════════════════

slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7)

add_text(slide7, "DEMO", 0.8, 0.3, 4, 0.8, size=42, color=CYAN, bold=True)
add_text(slide7, "Architektura demonstracyjna", 2.8, 0.45, 6, 0.5, size=20, color=GRAY)
shape = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2),
                                Inches(2), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = CYAN
shape.line.fill.background()

# ── Row 1: Browser + Web UI ──
add_labeled_box(slide7, "Przeglądarka", 0.5, 1.6, 2.0, 0.8, BG_CARD, PURPLE, WHITE, 14)
add_text(slide7, "PKCE Login", 0.6, 2.45, 1.8, 0.3, size=10, color=GRAY, align=PP_ALIGN.CENTER)

add_arrow(slide7, 2.6, 1.85, 3.3, 1.85, PURPLE, Pt(2))

# Agent Web UI
web_box = add_box(slide7, 3.4, 1.4, 2.8, 1.4, BG_CARD, CYAN)
add_text(slide7, "Agent Web UI", 3.5, 1.45, 2.6, 0.4, size=14, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide7, "Google ADK Agent", 3.5, 1.85, 2.6, 0.3, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide7, "FastAPI + Session", 3.5, 2.15, 2.6, 0.3, size=10, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide7, ":8000", 6.0, 2.5, 0.6, 0.3, size=10, color=GRAY, font_name="Consolas")

# ── Keycloak (center-top) ──
kc_box = add_box(slide7, 7.2, 1.4, 2.5, 1.6, BG_CARD, YELLOW)
add_text(slide7, "Keycloak", 7.3, 1.5, 2.3, 0.4, size=16, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)
add_text(slide7, "OIDC Provider", 7.3, 1.95, 2.3, 0.3, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide7, "Token Exchange", 7.3, 2.25, 2.3, 0.3, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide7, "JWKS / Realm Roles", 7.3, 2.55, 2.3, 0.3, size=10, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide7, ":8080", 9.5, 2.75, 0.6, 0.3, size=10, color=GRAY, font_name="Consolas")

# Arrow Web UI → Keycloak
add_arrow(slide7, 6.3, 1.85, 7.1, 1.85, YELLOW, Pt(2))

# ── Clients ──
add_text(slide7, "Keycloak Clients:", 10.2, 1.5, 2.8, 0.3, size=11, color=GRAY, bold=True)
clients = [
    ("adk-web-client", "public (PKCE)"),
    ("mcp-gateway", "confidential"),
    ("mcp-weather", "confidential"),
    ("mcp-calculator", "confidential"),
]
for ci, (cname, ctype) in enumerate(clients):
    add_text(slide7, f"▸ {cname}", 10.2, 1.85 + ci * 0.35, 1.8, 0.3,
             size=9, color=CYAN, font_name="Consolas")
    add_text(slide7, ctype, 11.9, 1.85 + ci * 0.35, 1.2, 0.3,
             size=9, color=GRAY)

# ── Row 2: Gateway ──
gw_demo_box = add_box(slide7, 3.4, 3.3, 5.5, 1.6, BG_CARD, GREEN)
add_text(slide7, "MCP Gateway", 3.5, 3.35, 2.5, 0.4, size=16, color=GREEN, bold=True)
add_text(slide7, ":8010", 8.5, 4.65, 0.6, 0.3, size=10, color=GRAY, font_name="Consolas")

# Gateway internals
gw_internals = [
    ("JWT Auth", 3.7, 3.85, CYAN, 1.5),
    ("Token Exchange", 5.4, 3.85, YELLOW, 1.7),
    ("search_servers", 3.7, 4.4, GREEN, 1.5),
    ("enable_server", 5.4, 4.4, GREEN, 1.7),
    ("Tool Proxy", 7.3, 4.1, PURPLE, 1.3),
]
for fname, fx, fy, fc, fw in gw_internals:
    feat = add_box(slide7, fx, fy, fw, 0.4, BG_CARD_LIGHT, fc)
    add_text(slide7, fname, fx + 0.05, fy + 0.02, fw - 0.1, 0.35,
             size=10, color=fc, bold=True, align=PP_ALIGN.CENTER, font_name="Consolas")

# Arrow Web → Gateway
add_line_arrow(slide7, 4.8, 2.8, 4.8, 3.3, CYAN)

# Arrow Gateway ↔ Keycloak
add_line_arrow(slide7, 8.45, 3.3, 8.45, 3.0, YELLOW)

# ── Row 3: MCP Servers ──
# Weather
weather_box = add_box(slide7, 2.5, 5.4, 3.0, 1.6, BG_CARD, TEAL)
add_text(slide7, "Weather MCP Server", 2.6, 5.45, 2.8, 0.35, size=13, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
add_text(slide7, "get_weather, get_forecast", 2.6, 5.85, 2.8, 0.3, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, font_name="Consolas")
add_text(slide7, "JWT: aud=mcp-weather", 2.6, 6.2, 2.8, 0.3, size=9, color=GRAY, align=PP_ALIGN.CENTER, font_name="Consolas")
add_text(slide7, "Open-Meteo API", 2.6, 6.5, 2.8, 0.3, size=10, color=TEAL, align=PP_ALIGN.CENTER)
add_text(slide7, ":8011", 5.3, 6.75, 0.6, 0.3, size=10, color=GRAY, font_name="Consolas")

# Calculator
calc_box = add_box(slide7, 6.5, 5.4, 3.0, 1.6, BG_CARD, TEAL)
add_text(slide7, "Calculator MCP Server", 6.6, 5.45, 2.8, 0.35, size=13, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
add_text(slide7, "add, subtract, multiply", 6.6, 5.85, 2.8, 0.3, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, font_name="Consolas")
add_text(slide7, "JWT: aud=mcp-calculator", 6.6, 6.2, 2.8, 0.3, size=9, color=GRAY, align=PP_ALIGN.CENTER, font_name="Consolas")
add_text(slide7, ":8012", 9.3, 6.75, 0.6, 0.3, size=10, color=GRAY, font_name="Consolas")

# Arrows Gateway → Servers
add_line_arrow(slide7, 5.0, 4.9, 4.0, 5.4, TEAL)
add_line_arrow(slide7, 7.0, 4.9, 8.0, 5.4, TEAL)

# Token flow summary at bottom
add_text(slide7, "Token flow:  User → PKCE login → token(aud:mcp-gateway) → Gateway → Token Exchange → token(aud:mcp-weather) → Server validates offline",
         0.3, 7.05, 12.5, 0.35, size=10, color=GRAY, align=PP_ALIGN.CENTER, font_name="Consolas")


# ════════════════════════════════════════════════════════════════
# SLIDE 8: Thank You
# ════════════════════════════════════════════════════════════════

slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8)

# Decorative circles
c1 = add_circle(slide8, -1, -1, 4, CYAN)
c1_sp = c1._element
solid = c1_sp.find(qn('p:spPr')).find(qn('a:solidFill'))
if solid is not None:
    srgb = solid.find(qn('a:srgbClr'))
    if srgb is not None:
        alpha = srgb.makeelement(qn('a:alpha'), {'val': '8000'})
        srgb.append(alpha)

c2 = add_circle(slide8, 10, 4, 5, PURPLE)
c2_sp = c2._element
solid = c2_sp.find(qn('p:spPr')).find(qn('a:solidFill'))
if solid is not None:
    srgb = solid.find(qn('a:srgbClr'))
    if srgb is not None:
        alpha = srgb.makeelement(qn('a:alpha'), {'val': '8000'})
        srgb.append(alpha)

add_text(slide8, "Dziękuję!", 1.5, 2.0, 10, 1.2, size=64, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

# Separator
shape = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(3.5),
                                Inches(3.3), Inches(0.05))
shape.fill.solid()
shape.fill.fore_color.rgb = CYAN
shape.line.fill.background()

add_text(slide8, "Adam Dąbrowski", 1.5, 3.8, 10, 0.6, size=24, color=WHITE, bold=False, align=PP_ALIGN.CENTER)

# Repo link
add_text(slide8, "Kod źródłowy / Demo:", 1.5, 4.8, 10, 0.5, size=16, color=GRAY, align=PP_ALIGN.CENTER)

# Repo URL in a styled box
repo_box = add_box(slide8, 3.5, 5.4, 6.3, 0.7, BG_CARD, CYAN)
add_text(slide8, "github.com/riposta/dynamic-mcp-poc", 3.7, 5.45, 5.9, 0.6,
         size=18, color=CYAN, bold=True, align=PP_ALIGN.CENTER, font_name="Consolas")

add_text(slide8, "Pytania?", 1.5, 6.4, 10, 0.6, size=28, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────

output_path = "./MCP_Gateway_Workshop.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
